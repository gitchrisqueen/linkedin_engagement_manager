import os
import random
import tempfile
import urllib.request
from datetime import datetime
from typing import Optional, Union

from lxml.etree import tostring
from pptx import Presentation
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.oxml import parse_xml
from pptx.slide import Slide
from pydantic import BaseModel, Field, HttpUrl, conlist, StrictStr
from pydantic_extra_types.color import Color


# Generic slide model for reusability
class CarouselSlide(BaseModel):
    title: Optional[StrictStr] = Field(None, description="Title or heading of the slide")
    content: Optional[StrictStr] = Field(None, description="Main content of the slide", max_length=500)
    image_url: Optional[HttpUrl] = Field(None, description="URL to an image")
    image_path: Optional[StrictStr] = Field(None, description="Path to an image")


# Specific Carousel Templates

class EducationalContentSlide(CarouselSlide):
    # Custom Educational Content slides could have additional elements if needed
    pass


class EducationalContentCarousel(BaseModel):
    cover: EducationalContentSlide = Field(..., description="Cover Slide: A bold title that clearly states the topic.")
    contents: conlist(EducationalContentSlide, min_length=1, max_length=4) = Field(...,
                                                                                   description="Content Slides: Each slide covers one tip or step with a combination of short text and relevant visuals.")
    call_to_action: EducationalContentSlide = Field(...,
                                                    description="Final Slide: Summarize key points and add a call to action.")


class CaseStudySlide(CarouselSlide):
    # Custom case study slides can have additional constraints or attributes
    pass


class CaseStudyCarousel(BaseModel):
    cover: CaseStudySlide = Field(...,
                                  description="Cover Slide: Title of the case study with the client’s name or project outcome.")
    challenge: CaseStudySlide = Field(..., description="Slide 2: Brief description of the problem faced by the client.")
    solution: CaseStudySlide = Field(..., description="Slide 3: Explanation of the approach or solution.")
    results: CaseStudySlide = Field(..., description="Slide 4: Highlight measurable results with data or visuals.")
    testimonial: Optional[CaseStudySlide] = Field(None,
                                                  description="Slide 5 (Optional): Include a client quote or feedback.")
    call_to_action: CaseStudySlide = Field(...,
                                           description="Final Slide: Encourage viewers to reach out for more information.")


class PersonalStorySlide(CarouselSlide):
    # Custom personal story slide if needed
    pass


class PersonalStoryCarousel(BaseModel):
    cover: PersonalStorySlide = Field(..., description="Cover Slide: Compelling title introducing the personal story.")
    story_slides: conlist(PersonalStorySlide, min_length=1, max_length=3) = Field(...,
                                                                                  description="Slides 2-4: Key moments in the journey.")
    takeaway: PersonalStorySlide = Field(..., description="Slide 5: Summary or lessons learned from the experience.")
    call_to_action: PersonalStorySlide = Field(...,
                                               description="Final Slide: Encourage viewers to share their own stories or connect with you for further discussion.")


class IndustryInsightSlide(CarouselSlide):
    # Custom industry insight slide for trends and insights
    pass


class IndustryInsightsCarousel(BaseModel):
    cover: IndustryInsightSlide = Field(...,
                                        description="Cover Slide: Title with an attention-grabbing phrase for industry insights.")
    insights: conlist(IndustryInsightSlide, min_length=1, max_length=4) = Field(...,
                                                                                description="Slides 2-5: Individual trends or insights with visuals.")
    call_to_action: IndustryInsightSlide = Field(...,
                                                 description="Final Slide: Summary and call-to-action for opinions on the trends.")


class EventRecapSlide(CarouselSlide):
    # Custom event recap slide if additional details are needed
    pass


class EventRecapCarousel(BaseModel):
    cover: EventRecapSlide = Field(..., description="cover Slide: Event title and date.")
    key_moments: conlist(EventRecapSlide, min_length=1, max_length=3) = Field(...,
                                                                              description="Slides 2-4: Key takeaways or highlights from the event.")
    call_to_action: EventRecapSlide = Field(...,
                                            description="Final Slide: Thank attendees and provide a call-to-action for future events or download additional resources.")


class TestimonialSlide(CarouselSlide):
    client_name: str = Field(..., description="Name of the client providing the testimonial")
    client_logo_url: Optional[HttpUrl] = Field(None, description="URL to the client’s logo or photo")


class TestimonialCarousel(BaseModel):
    cover: CarouselSlide = Field(...,
                                 description="Cover Slide: Cover slide with a title like 'What Our Clients Are Saying'.")
    testimonials: conlist(TestimonialSlide, min_length=1, max_length=3) = Field(...,
                                                                                description="Slides 2-4: Individual testimonials with a quote, client name, and photo.")
    call_to_action: CarouselSlide = Field(...,
                                          description="Final Slide: Call-to-action to encourage viewers to reach out for similar results.")


class ProductDemoSlide(CarouselSlide):
    # Custom product demo slide if additional details are needed
    pass


class ProductDemoCarousel(BaseModel):
    cover: ProductDemoSlide = Field(...,
                                    description="Cover Slide: Introduction to the product with a compelling headline.")
    main_feature: ProductDemoSlide = Field(...,
                                           description="Slide 2: Highlight the main feature of the product with an image.")
    additional_features: conlist(ProductDemoSlide, min_length=1, max_length=2) = Field(...,
                                                                                       description="Slides 3-4: Additional features of the product.")
    call_to_action: ProductDemoSlide = Field(...,
                                             description="Final Slide: Call-to-action to learn more or sign up for a demo.")


class PowerPointThemeColors(BaseModel):
    dk1: Optional[Color] = Field(None, description="RGB color code for the 1st dark color in the theme")
    lt1: Optional[Color] = Field(None, description="RGB color code for the 1st light color in the theme")
    dk2: Optional[Color] = Field(None, description="RGB color code for the 2nd dark color in the theme")
    lt2: Optional[Color] = Field(None, description="RGB color code for the 2nd light color in the theme")
    accent1: Optional[Color] = Field(None, description="RGB color code for the 1st accent color in the theme")
    accent2: Optional[Color] = Field(None, description="RGB color code for the 2nd accent color in the theme")
    accent3: Optional[Color] = Field(None, description="RGB color code for the 3rd accent color in the theme")
    accent4: Optional[Color] = Field(None, description="RGB color code for the 4th accent color in the theme")
    accent5: Optional[Color] = Field(None, description="RGB color code for the 5th accent color in the theme")
    accent6: Optional[Color] = Field(None, description="RGB color code for the 6th accent color in the theme")
    hlink: Optional[Color] = Field(None, description="RGB color code for the hyperlink color in the theme")
    folHlink: Optional[Color] = Field(None,
                                      description="RGB color code for the followed hyperlink color in the theme")


def create_ppt(ppt_name, carousel_data: Union[
    EducationalContentCarousel,
    CaseStudyCarousel,
    PersonalStoryCarousel,
    IndustryInsightsCarousel,
    EventRecapCarousel,
    TestimonialCarousel,
    ProductDemoCarousel],
               my_theme: PowerPointThemeColors = PowerPointThemeColors(**{"lt1": "e9d437", "dk2": "a89816"}),
               design_number: int = 1):
    current_dir = os.path.dirname(__file__)
    generated_dir = os.path.join(current_dir, "generated_designs")
    os.makedirs(generated_dir, exist_ok=True)
    design_path = os.path.join(current_dir, f"carousel_designs/Design-{design_number}.pptx")
    prs = Presentation(design_path)

    if isinstance(carousel_data, EducationalContentCarousel):
        # Handle EducationalContentCarousel
        prs = create_ppt_educational_content_carousel(prs, carousel_data)
        pass
    elif isinstance(carousel_data, CaseStudyCarousel):
        # Handle CaseStudyCarousel
        prs = create_ppt_case_study_carousel(prs, carousel_data)
        pass
    elif isinstance(carousel_data, PersonalStoryCarousel):
        prs = create_ppt_personal_story_carousel(prs, carousel_data)
    elif isinstance(carousel_data, IndustryInsightsCarousel):
        prs = create_ppt_industry_insights_carousel(prs, carousel_data)
    elif isinstance(carousel_data, EventRecapCarousel):
        prs = create_ppt_event_recap_carousel(prs, carousel_data)
    elif isinstance(carousel_data, TestimonialCarousel):
        prs = create_ppt_testimonial_carousel(prs, carousel_data)
    elif isinstance(carousel_data, ProductDemoCarousel):
        prs = create_ppt_product_demo_carousel(prs, carousel_data)

    file_path = os.path.join(generated_dir, f"{ppt_name}.pptx")
    prs.save(file_path)

    convert_ppt_theme_colors(file_path, my_theme)

    return f"{file_path}"


def get_default_image_path() -> str:
    # Get the default image path local to this file
    file_dir = os.path.dirname(__file__)
    default_image_path = os.path.join(file_dir, "images/image.png")
    return default_image_path


def get_pexels_image_path(query: str, default_path: Optional[str] = None) -> Optional[str]:
    """Download a Pexels image matching *query* to a temp file and return its path.

    Falls back to *default_path* when PEXELS_API_KEY is absent or the request fails.
    Pass default_path=None to get None back on failure (callers that must NOT post a
    placeholder — e.g. the carousel poster — rely on this to flag the post instead).
    """
    try:
        from cqc_lem.utilities.pexels_helper import get_photo
        photo = get_photo(query)
        url = photo.medium  # medium-size JPEG is a good balance for slides
        suffix = ".jpg"
        tmp = tempfile.NamedTemporaryFile(delete=False, suffix=suffix)
        urllib.request.urlretrieve(url, tmp.name)
        return tmp.name
    except Exception:
        return default_path


def create_ppt_educational_content_carousel(prs: Presentation, carousel: EducationalContentCarousel) -> Presentation:
    """
    Create a PowerPoint presentation for Educational Content Carousel.

    Parameters:
    - prs: Presentation object to add slides to.
    - carousel: EducationalContentCarousel containing carouseldata.

    """

    # Get the default image path local to this file
    default_image_path = get_default_image_path()

    # Slide 1: Cover
    cover_layouts = [create_title_layout_slide, create_title_only_layout_slide,
                     create_title_and_body_layout_slide]
    cover_slide_args = {
        "prs": prs,
        "title": carousel.cover.title,
        "subtitle": carousel.cover.content,
        "body_text": carousel.cover.content
    }
    cover_slide = random.choice(cover_layouts)(
        **cover_slide_args
    )

    # Slide 2-5: Content/Tips
    content_layouts = [create_title_and_body_layout_slide, create_one_column_text_layout_slide,
                       # create_ppt_for_title_and_two_columns_layout # TODO: Figure out if and how to implement this one
                       ]

    content_slide_args = {
        "prs": prs
    }
    for content in carousel.contents:
        content_slide_args["title"] = content.title
        content_slide_args["body_text"] = content.content
        image_path = getattr(content, "image_path", None)
        content_slide_args["image_path"] = image_path or get_pexels_image_path(
            content.title or "professional", default_image_path
        )
        content_slide = random.choice(content_layouts)(
            **content_slide_args
        )

    # Slide 6: Conclusion
    conclusion_layouts = [create_section_title_and_description_layout_slide, create_custom_3_1_layout_slide,
                          create_caption_only_layout_slide]
    cta_image_path = getattr(carousel.call_to_action, "image_path", None) or get_pexels_image_path(
        carousel.call_to_action.title or "success", default_image_path
    )
    conclusion_slide_args = {
        "prs": prs,
        "title": carousel.call_to_action.title,
        "description": carousel.call_to_action.content,
        "subtitle": carousel.call_to_action.content,
        "image_path": cta_image_path,
    }
    conclusion_slide = random.choice(conclusion_layouts)(
        **conclusion_slide_args
    )

    return prs  # Return the presentation for further modifications or saving


def create_ppt_case_study_carousel(prs: Presentation, case_study_carousel: CaseStudyCarousel) -> Presentation:
    """
    Create a PowerPoint presentation for Case Study Carousel.

    Parameters:
    - prs: Presentation object to add slides to.
    - case_study_carousel: CaseStudyCarousel instance with structured content for each slide.
    """

    # Slide 1: Cover
    cover_layouts = [create_title_layout_slide, create_section_header_layout_slide, create_title_only_layout_slide]
    cover_kwargs = {
        'prs': prs,
        'percentage': getattr(case_study_carousel.cover, 'percentage', ''),
        'title': case_study_carousel.cover.title,
        'subtitle': getattr(case_study_carousel.cover, 'subtitle', case_study_carousel.cover.content)
    }
    cover_slide = random.choice(cover_layouts)(**cover_kwargs)

    # Slide 2: Challenge
    challenge_layouts = [create_title_and_body_layout_slide, create_one_column_text_layout_slide]
    challenge_kwargs = {
        'prs': prs,
        'title': case_study_carousel.challenge.title,
        'body_text': case_study_carousel.challenge.content,
        'image_path': getattr(case_study_carousel.challenge, 'image_path', None) or get_pexels_image_path(
            case_study_carousel.challenge.title or "challenge", get_default_image_path()
        ),
    }
    challenge_slide = random.choice(challenge_layouts)(**challenge_kwargs)

    # Slide 3: Solution
    solution_layouts = [create_title_and_body_layout_slide, create_title_and_body_1_layout_slide,
                        # create_title_and_two_columns_layout_slide # TODO: Figure how to implement this one
                        ]
    solution_kwargs = {
        'prs': prs,
        'title': case_study_carousel.solution.title,
        'body_text': case_study_carousel.solution.content
    }
    solution_slide = random.choice(solution_layouts)(**solution_kwargs)

    # Slide 4: Results
    results_layouts = [create_big_number_layout_slide, create_title_and_body_layout_slide,
                       create_one_column_text_layout_slide]
    results_kwargs = {
        'prs': prs,
        'title': case_study_carousel.results.title,
        'body_text': case_study_carousel.results.content,
        'image_path': getattr(case_study_carousel.results, 'image_path', None) or get_pexels_image_path(
            case_study_carousel.results.title or "results", get_default_image_path()
        ),
        'big_number': getattr(case_study_carousel.results, 'big_number', ''),
        'subtitle': getattr(case_study_carousel.results, 'subtitle', case_study_carousel.results.content),
    }
    results_slide = random.choice(results_layouts)(**results_kwargs)

    # Slide 5: Testimonial
    testimonial_layouts = [create_caption_only_layout_slide, create_blank_1_1_layout_slide]
    testimonial_kwargs = {
        'prs': prs,
        'image_path': getattr(case_study_carousel.testimonial, 'image_path', get_default_image_path()),
        'title': getattr(case_study_carousel.testimonial, 'title', "Testimonial"),
        'quote': getattr(case_study_carousel.testimonial, 'quote', case_study_carousel.testimonial.content),
        'author': getattr(case_study_carousel.testimonial, 'author',
                          getattr(case_study_carousel.testimonial, 'title', "Happy Client"))
    }
    testimonial_slide = random.choice(testimonial_layouts)(**testimonial_kwargs)

    # Slide 6: Call to Action
    cta_layouts = [create_section_title_and_description_layout_slide, create_custom_3_1_layout_slide]
    cta_kwargs = {
        'prs': prs,
        'title': case_study_carousel.call_to_action.title,
        'subtitle': case_study_carousel.call_to_action.content,
        'description': case_study_carousel.call_to_action.content
    }
    cta_slide = random.choice(cta_layouts)(**cta_kwargs)

    return prs  # Return the presentation for further modifications or saving


def create_ppt_personal_story_carousel(prs: Presentation, carousel: PersonalStoryCarousel) -> Presentation:
    default_image = get_default_image_path()

    cover_layouts = [create_title_layout_slide, create_section_header_layout_slide, create_title_only_layout_slide]
    random.choice(cover_layouts)(
        prs=prs, title=carousel.cover.title, subtitle=carousel.cover.content,
        percentage="", body_text=carousel.cover.content
    )

    story_layouts = [create_title_and_body_layout_slide, create_one_column_text_layout_slide]
    for slide_data in carousel.story_slides:
        random.choice(story_layouts)(
            prs=prs, title=slide_data.title, body_text=slide_data.content,
            image_path=getattr(slide_data, "image_path", None) or get_pexels_image_path(
                slide_data.title or "story", default_image
            ),
        )

    create_title_and_body_1_layout_slide(
        prs=prs, title=carousel.takeaway.title, body_text=carousel.takeaway.content
    )

    cta_layouts = [create_section_title_and_description_layout_slide, create_custom_3_1_layout_slide]
    random.choice(cta_layouts)(
        prs=prs, title=carousel.call_to_action.title,
        description=carousel.call_to_action.content, subtitle=carousel.call_to_action.content
    )
    return prs


def create_ppt_industry_insights_carousel(prs: Presentation, carousel: IndustryInsightsCarousel) -> Presentation:
    default_image = get_default_image_path()

    cover_layouts = [create_title_layout_slide, create_title_only_layout_slide]
    random.choice(cover_layouts)(
        prs=prs, title=carousel.cover.title, subtitle=carousel.cover.content
    )

    insight_layouts = [create_title_and_body_layout_slide, create_one_column_text_layout_slide]
    for insight in carousel.insights:
        random.choice(insight_layouts)(
            prs=prs, title=insight.title, body_text=insight.content,
            image_path=getattr(insight, "image_path", None) or get_pexels_image_path(
                insight.title or "industry", default_image
            ),
        )

    cta_layouts = [create_section_title_and_description_layout_slide, create_custom_3_1_layout_slide]
    random.choice(cta_layouts)(
        prs=prs, title=carousel.call_to_action.title,
        description=carousel.call_to_action.content, subtitle=carousel.call_to_action.content
    )
    return prs


def create_ppt_event_recap_carousel(prs: Presentation, carousel: EventRecapCarousel) -> Presentation:
    default_image = get_default_image_path()

    cover_layouts = [create_title_layout_slide, create_section_header_layout_slide]
    random.choice(cover_layouts)(
        prs=prs, title=carousel.cover.title, subtitle=carousel.cover.content,
        percentage=""
    )

    moment_layouts = [create_title_and_body_layout_slide, create_one_column_text_layout_slide]
    for moment in carousel.key_moments:
        random.choice(moment_layouts)(
            prs=prs, title=moment.title, body_text=moment.content,
            image_path=getattr(moment, "image_path", None) or get_pexels_image_path(
                moment.title or "event", default_image
            ),
        )

    cta_layouts = [create_section_title_and_description_layout_slide, create_custom_3_1_layout_slide]
    random.choice(cta_layouts)(
        prs=prs, title=carousel.call_to_action.title,
        description=carousel.call_to_action.content, subtitle=carousel.call_to_action.content
    )
    return prs


def create_ppt_testimonial_carousel(prs: Presentation, carousel: TestimonialCarousel) -> Presentation:
    cover_layouts = [create_title_layout_slide, create_title_only_layout_slide]
    random.choice(cover_layouts)(
        prs=prs, title=carousel.cover.title, subtitle=carousel.cover.content
    )

    for testimonial in carousel.testimonials:
        create_blank_1_1_layout_slide(
            prs=prs,
            quote=f'"{testimonial.content}"',
            author=f"— {testimonial.client_name}"
        )

    cta_layouts = [create_section_title_and_description_layout_slide, create_custom_3_1_layout_slide]
    random.choice(cta_layouts)(
        prs=prs, title=carousel.call_to_action.title,
        description=carousel.call_to_action.content, subtitle=carousel.call_to_action.content
    )
    return prs


def create_ppt_product_demo_carousel(prs: Presentation, carousel: ProductDemoCarousel) -> Presentation:
    default_image = get_default_image_path()

    cover_layouts = [create_title_layout_slide, create_title_only_layout_slide]
    random.choice(cover_layouts)(
        prs=prs, title=carousel.cover.title, subtitle=carousel.cover.content
    )

    feature_layouts = [create_title_and_body_layout_slide, create_one_column_text_layout_slide]
    random.choice(feature_layouts)(
        prs=prs, title=carousel.main_feature.title, body_text=carousel.main_feature.content,
        image_path=getattr(carousel.main_feature, "image_path", None) or get_pexels_image_path(
            carousel.main_feature.title or "product", default_image
        ),
    )

    for feature in carousel.additional_features:
        random.choice(feature_layouts)(
            prs=prs, title=feature.title, body_text=feature.content,
            image_path=getattr(feature, "image_path", None) or get_pexels_image_path(
                feature.title or "feature", default_image
            ),
        )

    cta_layouts = [create_section_title_and_description_layout_slide, create_custom_3_1_layout_slide]
    random.choice(cta_layouts)(
        prs=prs, title=carousel.call_to_action.title,
        description=carousel.call_to_action.content, subtitle=carousel.call_to_action.content
    )
    return prs


def debug_slide(slide):
    for shape in slide.shapes:
        print(f"Shape: {shape.name}, Type: {shape.shape_type}, Placeholder: {shape.placeholder_format.idx}")


def convert_ppt_theme_colors(ppt_path, theme_colors: PowerPointThemeColors):
    # Load the presentation
    prs = Presentation(ppt_path)

    # Get the Slide Master
    slide_master = prs.slide_master
    slide_master_part = slide_master.part

    # Get the Theme and part
    theme_part = slide_master_part.part_related_by(RT.THEME)
    theme = parse_xml(theme_part.blob)  # theme here is an <a:theme> element

    # For each of the attributes in the PowerPointThemeColors model, find the corresponding XML element and update the color value
    for field_name, field_value in theme_colors:
        if field_value:
            color_element = theme.xpath(f'a:themeElements/a:clrScheme/a:{field_name}/a:srgbClr')[0]
            # print(f"{field_name} color before: {color_element.get('val')}")
            set_color = field_value.as_hex(format="long").replace("#", "")
            # print(f"{field_name} color after: {set_color}")
            color_element.set('val', set_color.encode('utf-8'))

    # Serialize the modified XML back to the theme part
    theme_part._blob = tostring(theme)

    # print(f"Blob After: {theme_part.blob}")

    # Save the presentation
    prs.save(ppt_path)


def set_ppt_theme_colors(ppt_path, theme_colors: dict = None):
    # Load default theme colors
    if theme_colors is None:
        theme_colors = {
            "dk1": "ffffff",
            "lt1": "ffffff",
            "dk2": "ffffff",
            "lt2": "ffffff",
            "accent1": "ffffff",
            "accent2": "ffffff",
            "accent3": "ffffff",
            "accent4": "ffffff",
            "accent5": "ffffff",
            "accent6": "ffffff",
            "hlink": "ffffff",
            "folHlink": "ffffff"
        }

    # Load the presentation
    prs = Presentation(ppt_path)

    # Get the Slide Master
    slide_master = prs.slide_master
    slide_master_part = slide_master.part

    # Get the Theme and part
    theme_part = slide_master_part.part_related_by(RT.THEME)
    theme = parse_xml(theme_part.blob)  # theme here is an <a:theme> element

    # For each of the theme color names in the themes dict, find the corresponding XML element and update the color value
    for theme_color_name, theme_color_hex_value in theme_colors.items():
        if theme_color_name:
            color_element = theme.xpath(f'a:themeElements/a:clrScheme/a:{theme_color_name}/a:srgbClr')[0]
            # print(f"{theme_color_name} color before: {color_element.get('val')}")
            color_element.set('val', theme_color_hex_value.encode('utf-8'))
            # print(f"{theme_color_name} color after: {theme_color_hex_value}")

    # Serialize the modified XML back to the theme part
    theme_part._blob = tostring(theme)

    # Save the presentation
    prs.save(ppt_path)


def get_attr_gracefully(obj, attr):
    try:
        return getattr(obj, attr)
    except Exception as e:
        print(f"Error accessing attribute ({attr}): {e}")
    return None


def create_title_layout_slide(prs: Presentation, title: str, subtitle: str, **kwargs) -> Slide:
    """Create a PowerPoint slide with the TITLE layout.

    Parameters:
    - prs: Presentation object to add slides to.
    - title: The main title text.
    - subtitle: Additional descriptive text.
    - **kwargs: Additional keyword arguments are thrown away to allow for flexible function calls.
    """

    # Get the layout and add a slide
    layout = prs.slide_layouts[0]  # Assuming 0 index is the TITLE layout
    slide = prs.slides.add_slide(layout)

    # Populate title and subtitle placeholders
    title_placeholder = slide.shapes.title
    subtitle_placeholder = slide.placeholders[1]  # Assuming index 1 is for subtitle

    # Set the text content
    title_placeholder.text = title
    subtitle_placeholder.text = subtitle

    return slide


def create_section_header_layout_slide(prs: Presentation, percentage: str, title: str, subtitle: str,
                                       **kwargs) -> Slide:
    """
    Create a PowerPoint slide with the SECTION_HEADER layout.

    Parameters:
    - prs: Presentation object to add slides to.
    - percentage: The main percentage or metric to highlight.
    - title: The main title text.
    - subtitle: Additional descriptive text.
    - **kwargs: Additional keyword arguments are thrown away to allow for flexible function calls.
    """
    # Get the layout and add a slide
    layout = prs.slide_layouts[1]  # Assuming 1 index is the SECTION_HEADER layout
    slide = prs.slides.add_slide(layout)

    # Populate the placeholders
    percentage_placeholder = slide.placeholders[0]
    title_placeholder = slide.placeholders[1]
    subtitle_placeholder = slide.placeholders[2]

    # Set the text content
    percentage_placeholder.text = percentage
    title_placeholder.text = title
    subtitle_placeholder.text = subtitle

    return slide


def debug_master_slide_placeholders_and_text(design_number: int = 1):
    current_dir = os.path.dirname(__file__)
    design_path = os.path.join(current_dir, f"carousel_designs/Design-{design_number}.pptx")
    prs = Presentation(design_path)
    slide_master = prs.slide_master
    for slide_layout in slide_master.slide_layouts:
        print(f"Slide Layout: {slide_layout.name}")
        for placeholder in slide_layout.placeholders:
            print(f"Placeholder: {placeholder.name}, Type: {placeholder.placeholder_format.idx}")
            for shape in slide_layout.shapes:
                if shape.is_placeholder:
                    phf = shape.placeholder_format
                    print(f"Shape: {shape.name}, Type: {shape.shape_type}, Placeholder: {phf.idx}")
                    if phf.idx == placeholder.placeholder_format.idx:
                        print(f"Text: {shape.text}")
        print("----")


def create_title_and_body_layout_slide(prs: Presentation, title: str, body_text: str, **kwargs) -> Slide:
    """
    Create a PowerPoint slide with the TITLE_AND_BODY layout.

    Parameters:
    - prs: Presentation object to add slides to.
    - title: The main title text for the slide.
    - body_text: The body text content, which could include bullet points or paragraphs.
    - **kwargs: Additional keyword arguments are thrown away to allow for flexible function calls.
    """
    # Get the layout and add a slide
    layout = prs.slide_layouts[2]  # Assuming 2 index is the TITLE_AND_BODY layout
    slide = prs.slides.add_slide(layout)

    # Populate the placeholders
    title_placeholder = slide.placeholders[0]
    body_placeholder = slide.placeholders[1]

    # Set the text content
    title_placeholder.text = title
    body_placeholder.text = body_text

    return slide


def create_title_and_two_columns_layout_slide(prs: Presentation, title: str,
                                              left_column_title: str, left_column_subtitle: str,
                                              right_column_title: str, right_column_subtitle: str,
                                              **kwargs
                                              ) -> Slide:
    """
    Create a PowerPoint slide with the TITLE_AND_TWO_COLUMNS layout.

    Parameters:
    - prs: Presentation object to add slides to.
    - title: The main title text for the slide.
    - left_column_title: Title for the left column.
    - left_column_subtitle: Subtitle text for the left column.
    - right_column_title: Title for the right column.
    - right_column_subtitle: Subtitle text for the right column.
    - **kwargs: Additional keyword arguments are thrown away to allow for flexible function calls.
    """
    # Get the layout and add a slide
    layout = prs.slide_layouts[3]  # Assuming 3 index is the TITLE_AND_TWO_COLUMNS layout
    slide = prs.slides.add_slide(layout)

    # Populate the placeholders
    title_placeholder = slide.placeholders[0]
    left_column_title_placeholder = slide.placeholders[1]
    left_column_subtitle_placeholder = slide.placeholders[2]
    right_column_title_placeholder = slide.placeholders[3]
    right_column_subtitle_placeholder = slide.placeholders[4]

    # Set the text content
    title_placeholder.text = title
    left_column_title_placeholder.text = left_column_title
    left_column_subtitle_placeholder.text = left_column_subtitle
    right_column_title_placeholder.text = right_column_title
    right_column_subtitle_placeholder.text = right_column_subtitle

    return slide


def create_title_only_layout_slide(prs: Presentation, title: str, **kwargs) -> Slide:
    """
    Create a PowerPoint slide with the TITLE_ONLY layout.

    Parameters:
    - prs: Presentation object to add slides to.
    - title: The main title text for the slide.
    - **kwargs: Additional keyword arguments are thrown away to allow for flexible function calls.

    Returns:
    - Slide: The PowerPoint slide object, allowing further customization.
    """
    # Get the layout and add a slide
    layout = prs.slide_layouts[4]  # Assuming 4 index is the TITLE_ONLY layout
    slide = prs.slides.add_slide(layout)

    # Populate the title placeholder
    title_placeholder = slide.placeholders[0]
    title_placeholder.text = title

    # Return the slide object for further customization
    return slide


def create_main_point_layout_slide(prs: Presentation, title: str, **kwargs) -> Slide:
    """
    Create a PowerPoint slide with the MAIN_POINT layout.

    Parameters:
    - prs: Presentation object to add slides to.
    - title: The main title text for the slide, typically a key point or highlight.
    - **kwargs: Additional keyword arguments are thrown away to allow for flexible function calls.

    Returns:
    - Slide: The PowerPoint slide object, allowing further customization.
    """
    # Get the layout and add a slide
    layout = prs.slide_layouts[6]  # Assuming 6 index is the MAIN_POINT layout
    slide = prs.slides.add_slide(layout)

    # Populate the title placeholder
    title_placeholder = slide.placeholders[0]
    title_placeholder.text = title

    # Return the slide object for further customization
    return slide


def create_one_column_text_layout_slide(prs: Presentation, title: str, body_text: str, image_path: str = None,
                                        **kwargs) -> Slide:
    """
    Create a PowerPoint slide with the ONE_COLUMN_TEXT layout.

    Parameters:
    - prs: Presentation object to add slides to.
    - title: The main title text for the slide.
    - body_text: The body text content for the left column.
    - image_path (optional): Path to the image file to be inserted into the picture placeholder.
    - **kwargs: Additional keyword arguments are thrown away to allow for flexible function calls.

    Returns:
    - Slide: The PowerPoint slide object, allowing further customization.
    """
    # Get the layout and add a slide
    layout = prs.slide_layouts[5]  # Assuming 5 index is the ONE_COLUMN_TEXT layout
    slide = prs.slides.add_slide(layout)

    # Populate the placeholders
    title_placeholder = slide.placeholders[0]
    body_placeholder = slide.placeholders[1]
    picture_placeholder = slide.placeholders[2]

    # Set the text content
    title_placeholder.text = title
    body_placeholder.text = body_text

    # Insert the image if the path is provided
    if image_path:
        picture_placeholder.insert_picture(image_path)

    # Return the slide object for further customization
    return slide


def create_section_title_and_description_layout_slide(prs: Presentation, title: str, description: str,
                                                      **kwargs) -> Slide:
    """
    Create a PowerPoint slide with the SECTION_TITLE_AND_DESCRIPTION layout.

    Parameters:
    - prs: Presentation object to add slides to.
    - title: The main title text for the slide.
    - description: Subtitle or description text for additional context.
    - **kwargs: Additional keyword arguments are thrown away to allow for flexible function calls.

    Returns:
    - Slide: The PowerPoint slide object, allowing further customization.
    """
    # Get the layout and add a slide
    layout = prs.slide_layouts[7]  # Assuming 7 index is the SECTION_TITLE_AND_DESCRIPTION layout
    slide = prs.slides.add_slide(layout)

    # Populate the placeholders
    title_placeholder = slide.placeholders[0]
    description_placeholder = slide.placeholders[1]

    # Set the text content
    title_placeholder.text = title
    description_placeholder.text = description

    # Return the slide object for further customization
    return slide


def create_caption_only_layout_slide(prs: Presentation, title: str, image_path: str = None, **kwargs) -> Slide:
    """
    Create a PowerPoint slide with the CAPTION_ONLY layout.

    Parameters:
    - prs: Presentation object to add slides to.
    - title: The caption or title text for the slide.
    - image_path: Path to the image file to be inserted into the picture placeholder.
    - **kwargs: Additional keyword arguments are thrown away to allow for flexible function calls.

    Returns:
    - Slide: The PowerPoint slide object, allowing further customization.
    """
    # Get the layout and add a slide
    layout = prs.slide_layouts[8]  # Assuming 8 index is the CAPTION_ONLY layout
    slide = prs.slides.add_slide(layout)

    # Populate the title placeholder
    title_placeholder = slide.placeholders[0]
    title_placeholder.text = title

    # Insert the image if the path is provided
    picture_placeholder = slide.placeholders[2]
    if image_path:
        picture_placeholder.insert_picture(image_path)

    # Return the slide object for further customization
    return slide


def create_big_number_layout_slide(prs: Presentation, big_number: str, subtitle: str, **kwargs) -> Slide:
    """
    Create a PowerPoint slide with the BIG_NUMBER layout.

    Parameters:
    - prs: Presentation object to add slides to.
    - big_number: The main number or percentage to highlight on the slide.
    - subtitle: Subtitle text for additional context or explanation.
    - **kwargs: Additional keyword arguments are thrown away to allow for flexible function calls.

    Returns:
    - Slide: The PowerPoint slide object, allowing further customization.
    """
    # Get the layout and add a slide
    layout = prs.slide_layouts[9]  # Assuming 9 index is the BIG_NUMBER layout
    slide = prs.slides.add_slide(layout)

    # Populate the placeholders
    big_number_placeholder = slide.placeholders[0]
    subtitle_placeholder = slide.placeholders[1]

    # Set the text content
    big_number_placeholder.text = big_number
    subtitle_placeholder.text = subtitle

    # Return the slide object for further customization
    return slide


def create_blank_layout_slide(prs: Presentation, **kwargs) -> Slide:
    """
    Create a PowerPoint slide with the BLANK layout.

    Parameters:
    - prs: Presentation object to add slides to.
    - **kwargs: Additional keyword arguments are thrown away to allow for flexible function calls.

    Returns:
    - Slide: The PowerPoint slide object, allowing complete customization.
    """
    # Get the layout and add a blank slide
    layout = prs.slide_layouts[10]  # Assuming 10 index is the BLANK layout
    slide = prs.slides.add_slide(layout)

    # Return the slide object for full customization
    return slide


def create_custom_6_1_layout_slide(prs: Presentation, title: str, columns: list[dict], **kwargs) -> Slide:
    """
    Create a PowerPoint slide with the CUSTOM_6_1 (Title and Three Columns) layout.

    Parameters:
    - prs: Presentation object to add slides to.
    - title: The main title text for the slide.
    - columns: A list of dictionaries, each containing 'metric', 'sub_title', and 'description' for each column.
               Example format: [{'metric': 'XX%', 'sub_title': 'Column 1 Title', 'description': 'Details for Column 1'}, ...]
    - **kwargs: Additional keyword arguments are thrown away to allow for flexible function calls.

    Returns:
    - Slide: The PowerPoint slide object, allowing further customization.
    """
    # Get the layout and add a slide
    layout = prs.slide_layouts[11]  # Assuming 11 index is the CUSTOM_6_1 layout
    slide = prs.slides.add_slide(layout)

    # Set the main title
    title_placeholder = slide.placeholders[0]
    title_placeholder.text = title

    # Populate each column's placeholders
    for i, column_content in enumerate(columns):
        if i > 2:
            break  # Only three columns available

        metric_placeholder = slide.placeholders[3 + i * 3]
        sub_title_placeholder = slide.placeholders[4 + i * 3]
        description_placeholder = slide.placeholders[5 + i * 3]

        # Set content for each column
        metric_placeholder.text = column_content['metric']
        sub_title_placeholder.text = column_content['sub_title']
        description_placeholder.text = column_content['description']

    # Return the slide object for further customization
    return slide


def create_title_only_1_1_layout_slide(prs: Presentation, title: str, **kwargs) -> Slide:
    """
    Create a PowerPoint slide with the TITLE_ONLY_1_1 layout.

    Parameters:
    - prs: Presentation object to add slides to.
    - title: The main title text for the slide.
    - **kwargs: Additional keyword arguments are thrown away to allow for flexible function calls.

    Returns:
    - Slide: The PowerPoint slide object, allowing further customization.
    """
    # Get the layout and add a slide
    layout = prs.slide_layouts[12]  # Assuming 12 index is the TITLE_ONLY_1_1 layout
    slide = prs.slides.add_slide(layout)

    # Populate the title placeholder
    title_placeholder = slide.placeholders[0]
    title_placeholder.text = title

    # Return the slide object for further customization
    return slide


def create_one_column_text_1_layout_slide(prs: Presentation, title: str, body_text: str, image_path: str = None,
                                          **kwargs) -> Slide:
    """
    Create a PowerPoint slide with the ONE_COLUMN_TEXT_1 layout.

    Parameters:
    - prs: Presentation object to add slides to.
    - title: The main title text for the slide.
    - body_text: The body text content for the right column.
    - image_path: Path to the image file to be inserted into the picture placeholder.
    - **kwargs: Additional keyword arguments are thrown away to allow for flexible function calls.

    Returns:
    - Slide: The PowerPoint slide object, allowing further customization.
    """
    # Get the layout and add a slide
    layout = prs.slide_layouts[13]  # Assuming 13 index is the ONE_COLUMN_TEXT_1 layout
    slide = prs.slides.add_slide(layout)

    # Populate the placeholders
    title_placeholder = slide.placeholders[0]
    picture_placeholder = slide.placeholders[1]
    body_placeholder = slide.placeholders[2]

    # Set the text content
    title_placeholder.text = title
    body_placeholder.text = body_text

    # Insert the image if the path is provided
    if image_path:
        picture_placeholder.insert_picture(image_path)

    # Return the slide object for further customization
    return slide


def create_blank_1_1_layout_slide(prs: Presentation, quote: str, author: str, **kwargs) -> Slide:
    """
    Create a PowerPoint slide with the BLANK_1_1 (Quote) layout.

    Parameters:
    - prs: Presentation object to add slides to.
    - quote: The main quote or message for the slide.
    - author: The author or source of the quote.
    - **kwargs: Additional keyword arguments are thrown away to allow for flexible function calls.

    Returns:
    - Slide: The PowerPoint slide object, allowing further customization.
    """
    # Get the layout and add a slide
    layout = prs.slide_layouts[14]  # Assuming 14 index is the BLANK_1_1 layout
    slide = prs.slides.add_slide(layout)

    # Populate the placeholders
    quote_placeholder = slide.placeholders[0]
    author_placeholder = slide.placeholders[1]

    # Set the text content
    quote_placeholder.text = quote
    author_placeholder.text = author

    # Return the slide object for further customization
    return slide


def create_title_and_two_columns_1_layout_slide(prs: Presentation, title: str, column_1: dict,
                                                column_2: dict, **kwargs) -> Slide:
    """
    Create a PowerPoint slide with the TITLE_AND_TWO_COLUMNS_1 layout.

    Parameters:
    - prs: Presentation object to add slides to.
    - title: The main title text for the slide.
    - column_1: A dictionary with 'sub_title' and 'text' for the left column.
               Example format: {'sub_title': 'Column 1 Title', 'text': 'Content for Column 1'}
    - column_2: A dictionary with 'sub_title' and 'text' for the right column.
               Example format: {'sub_title': 'Column 2 Title', 'text': 'Content for Column 2'}
    - **kwargs: Additional keyword arguments are thrown away to allow for flexible function calls.

    Returns:
    - Slide: The PowerPoint slide object, allowing further customization.
    """
    # Get the layout and add a slide
    layout = prs.slide_layouts[15]  # Assuming 15 index is the TITLE_AND_TWO_COLUMNS_1 layout
    slide = prs.slides.add_slide(layout)

    # Set the main title
    title_placeholder = slide.placeholders[0]
    title_placeholder.text = title

    # Populate left column's placeholders
    column_1_subtitle_placeholder = slide.placeholders[1]
    column_1_text_placeholder = slide.placeholders[2]
    column_1_subtitle_placeholder.text = column_1['sub_title']
    column_1_text_placeholder.text = column_1['text']

    # Populate right column's placeholders
    column_2_subtitle_placeholder = slide.placeholders[3]
    column_2_text_placeholder = slide.placeholders[4]
    column_2_subtitle_placeholder.text = column_2['sub_title']
    column_2_text_placeholder.text = column_2['text']

    # Return the slide object for further customization
    return slide


def create_title_and_body_1_layout_slide(prs: Presentation, title: str, body_text: str, **kwargs) -> Slide:
    """
    Create a PowerPoint slide with the TITLE_AND_BODY_1 layout.

    Parameters:
    - prs: Presentation object to add slides to.
    - title: The main title text for the slide.
    - body_text: The body text content for the slide, supporting paragraphs or bullet points.
    - **kwargs: Additional keyword arguments are thrown away to allow for flexible function calls.

    Returns:
    - Slide: The PowerPoint slide object, allowing further customization.
    """
    # Get the layout and add a slide
    layout = prs.slide_layouts[16]  # Assuming 16 index is the TITLE_AND_BODY_1 layout
    slide = prs.slides.add_slide(layout)

    # Populate the placeholders
    title_placeholder = slide.placeholders[0]
    body_placeholder = slide.placeholders[1]

    # Set the text content
    title_placeholder.text = title
    body_placeholder.text = body_text

    # Return the slide object for further customization
    return slide


def create_custom_3_1_layout_slide(prs: Presentation, title: str, subtitle: str, **kwargs) -> Slide:
    """
    Create a PowerPoint slide with the CUSTOM_3_1 (Thanks) layout.

    Parameters:
    - prs: Presentation object to add slides to.
    - title: The main title text for the slide, typically a closing or thank-you message.
    - subtitle: Additional information or call-to-action text.
    - **kwargs: Additional keyword arguments are thrown away to allow for flexible function calls.

    Returns:
    - Slide: The PowerPoint slide object, allowing further customization.
    """
    # Get the layout and add a slide
    layout = prs.slide_layouts[17]  # Assuming 17 index is the CUSTOM_3_1 (Thanks) layout
    slide = prs.slides.add_slide(layout)

    # Populate the placeholders
    title_placeholder = slide.placeholders[0]
    subtitle_placeholder = slide.placeholders[1]

    # Set the text content
    title_placeholder.text = title
    subtitle_placeholder.text = subtitle

    # Return the slide object for further customization
    return slide


def test_create_educational_ppt():
    """
    Test function to create an EducationalContentCarousel presentation and test the create_ppt function.
    """
    # Sample data for an educational content carousel
    carousel_data = {
        "cover": {
            "title": "5 Tips for Boosting Productivity",
            "content": "Practical ways to get more done every day"
        },
        "contents": [
            {"title": "Tip 1", "content": "Set clear goals for each day."},
            {"title": "Tip 2", "content": "Take regular breaks to recharge."},
            {"title": "Tip 3", "content": "Eliminate distractions to stay focused."},
            {"title": "Tip 4", "content": "Use productivity tools to track your progress."}
        ],
        "call_to_action": {
            "title": "Comment Below",
            "content": "Which Tip Will You Try First?"
        }
    }

    # PPT Name with tiemstamp suffix
    ppt_name = f"Educational_Carousel_{datetime.now().strftime('%Y%m%d_%H%M%S')}"
    my_ppt = create_ppt(ppt_name, EducationalContentCarousel(**carousel_data))
    print(f"Presentation created: {my_ppt}")


def test_create_case_study_ppt():
    """
    Test function to create an CaseStudyCarousel presentation and test the create_ppt function.
    """
    # Sample data for an case study content carousel
    carousel_data = {
        "cover": {
            "title": "Case Study: Successful Project",
            "content": "An in-depth look at our successful project with Client X"
        },
        "challenge": {
            "title": "The Challenge",
            "content": "Client X faced significant challenges in their market due to increased competition and changing customer preferences.",
            "image_path": get_default_image_path()
        },
        "solution": {
            "title": "Our Solution",
            "content": "We implemented a comprehensive strategy that included market analysis, customer engagement, and product innovation."
        },
        "results": {
            "title": "The Results",
            "content": "Our solution led to a 30% increase in market share and a 20% increase in customer satisfaction.",
            "image_path": get_default_image_path(),
            "big_number": "30%",
            "subtitle": "Increase in Market Share"
        },
        "testimonial": {
            "title": "Client Testimonial",
            "content": "Working with this team was a game-changer for our business. Their expertise and dedication were evident in every step of the process.",
            "image_path": get_default_image_path(),
            "quote": "This team transformed our business.",
            "author": "Jane Doe, CEO of Client X"
        },
        "call_to_action": {
            "title": "Get in Touch",
            "content": "Contact us to learn how we can help your business achieve similar results."
        }
    }

    # PPT Name with timestamp suffix
    ppt_name = f"Case_Study_Carousel_{datetime.now().strftime('%Y%m%d_%H%M%S')}"
    my_ppt = create_ppt(ppt_name, CaseStudyCarousel(**carousel_data))
    print(f"Presentation created: {my_ppt}")


def test_caption_only_slide(design_number: int = 1):
    current_dir = os.path.dirname(__file__)
    generated_dir = os.path.join(current_dir, "generated_designs")
    os.makedirs(generated_dir, exist_ok=True)
    design_path = os.path.join(current_dir, f"carousel_designs/Design-{design_number}.pptx")
    prs = Presentation(design_path)
    create_caption_only_layout_slide(prs, "Caption Only Slide", image_path=get_default_image_path())
    ppt_name = f"Caption Only Slide_{datetime.now().strftime('%Y%m%d_%H%M%S')}"
    file_path = os.path.join(generated_dir, f"{ppt_name}.pptx")
    prs.save(file_path)

    print(f"Created: {file_path}")


# ── Research-backed carousel templates ────────────────────────────────────────
# Each template is a dict of color and style parameters.
# Derived from Buffer/PostNitro/Hootsuite analysis of highest-engagement carousels.
CAROUSEL_TEMPLATES: dict[str, dict] = {
    "bold_listicle": {
        "label": "Bold Listicle",
        "description": "White slides, numbered badge circles, rainbow accents. Best for: tips, tools, mistakes.",
        "layout":      "listicle",
        "cover_bg":    (15, 23, 42),
        "cover_text":  (255, 255, 255),
        "cover_accent": (59, 130, 246),
        "content_bg":  (255, 255, 255),
        "title_color": (15, 23, 42),
        "body_color":  (71, 85, 105),
        "bottom_bar":  (15, 23, 42),
        "badge_colors": [(59, 130, 246), (16, 185, 129), (239, 68, 68), (139, 92, 246), (245, 158, 11), (14, 165, 233)],
    },
    "minimal_dark": {
        "label": "Minimal Dark",
        "description": "Black slides, huge left-aligned titles, gold accents. Best for: bold opinions, predictions.",
        "layout":      "dark_minimal",
        "cover_bg":    (10, 10, 10),
        "cover_text":  (255, 255, 255),
        "cover_accent": (251, 191, 36),
        "content_bg":  (18, 18, 18),
        "title_color": (255, 255, 255),
        "body_color":  (163, 163, 163),
        "bottom_bar":  (30, 30, 30),
        "badge_colors": [(251, 191, 36), (251, 146, 60), (52, 211, 153), (129, 140, 248), (248, 113, 113), (34, 211, 238)],
    },
    "stat_reveal": {
        "label": "Stat Reveal",
        "description": "Each slide title displayed ENORMOUS centered. Best for: data insights, research findings.",
        "layout":      "stat_big",
        "cover_bg":    (30, 64, 175),
        "cover_text":  (255, 255, 255),
        "cover_accent": (147, 197, 253),
        "content_bg":  (239, 246, 255),
        "title_color": (30, 64, 175),
        "body_color":  (55, 65, 81),
        "bottom_bar":  (30, 64, 175),
        "badge_colors": [(30, 64, 175)] * 6,
    },
    "step_framework": {
        "label": "Step Framework",
        "description": "Visual progress dots at top, arrow-bulleted body. Best for: how-to guides, playbooks.",
        "layout":      "step_progress",
        "cover_bg":    (4, 120, 87),
        "cover_text":  (255, 255, 255),
        "cover_accent": (110, 231, 183),
        "content_bg":  (255, 255, 255),
        "title_color": (6, 78, 59),
        "body_color":  (55, 65, 81),
        "bottom_bar":  (4, 120, 87),
        "badge_colors": [(16, 185, 129)] * 6,
    },
    "story_arc": {
        "label": "Story Arc",
        "description": "Cream slides, giant quote marks, square numbered badges. Best for: personal stories.",
        "layout":      "quote_pull",
        "cover_bg":    (120, 53, 15),
        "cover_text":  (255, 255, 255),
        "cover_accent": (253, 186, 116),
        "content_bg":  (255, 251, 235),
        "title_color": (92, 45, 0),
        "body_color":  (120, 53, 15),
        "bottom_bar":  (120, 53, 15),
        "badge_colors": [(245, 158, 11), (234, 88, 12), (217, 70, 239), (99, 102, 241), (239, 68, 68), (16, 185, 129)],
    },
}

DEFAULT_TEMPLATE = "bold_listicle"


def create_carousel_slide_images(
    carousel_data: Union[
        EducationalContentCarousel,
        CaseStudyCarousel,
        PersonalStoryCarousel,
        IndustryInsightsCarousel,
        EventRecapCarousel,
        TestimonialCarousel,
        ProductDemoCarousel,
    ],
    post_id: int,
    output_dir: Optional[str] = None,
    template: str = DEFAULT_TEMPLATE,
    # Legacy params retained for backward compat but ignored
    bg_color: tuple = (26, 86, 219),
    accent_color: tuple = (255, 255, 255),
    secondary_bg: tuple = (15, 52, 142),
) -> list[str]:
    """Render carousel slides as 1080x1080 PNG images using Pillow.

    Creates one image per slide in output_dir (defaults to
    assets/images/carousel/{post_id}/). Returns a list of absolute image paths.

    ``template`` selects a visual style from CAROUSEL_TEMPLATES. Defaults to
    DEFAULT_TEMPLATE ("bold_listicle").
    """
    from PIL import Image, ImageDraw, ImageFont

    W, H = 1080, 1080
    WHITE = (255, 255, 255)

    tmpl = CAROUSEL_TEMPLATES.get(template, CAROUSEL_TEMPLATES[DEFAULT_TEMPLATE])
    cover_bg     = tmpl["cover_bg"]
    cover_text   = tmpl["cover_text"]
    cover_accent = tmpl["cover_accent"]
    content_bg   = tmpl["content_bg"]
    title_color  = tmpl["title_color"]
    body_color   = tmpl["body_color"]
    bottom_bar   = tmpl["bottom_bar"]
    badge_colors = tmpl["badge_colors"]
    layout       = tmpl.get("layout", "listicle")

    if output_dir is None:
        current_dir = os.path.dirname(__file__)
        assets_root = os.path.join(current_dir, "..", "assets", "images", "carousel", str(post_id))
        output_dir = os.path.realpath(assets_root)
    os.makedirs(output_dir, exist_ok=True)

    # ── Font loader ───────────────────────────────────────────────────────────
    def _load_font(size: int, bold: bool = True) -> ImageFont.FreeTypeFont:
        bold_paths = [
            "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf",
            "/usr/share/fonts/truetype/liberation/LiberationSans-Bold.ttf",
            "/usr/share/fonts/truetype/freefont/FreeSansBold.ttf",
            "/System/Library/Fonts/Supplemental/Arial Bold.ttf",
            "/Library/Fonts/Arial Bold.ttf",
            "/System/Library/Fonts/Helvetica.ttc",
        ]
        reg_paths = [
            "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
            "/usr/share/fonts/truetype/liberation/LiberationSans-Regular.ttf",
            "/usr/share/fonts/truetype/freefont/FreeSans.ttf",
            "/System/Library/Fonts/Supplemental/Arial.ttf",
            "/Library/Fonts/Arial.ttf",
            "/System/Library/Fonts/Helvetica.ttc",
        ]
        for path in (bold_paths if bold else reg_paths):
            if os.path.exists(path):
                try:
                    return ImageFont.truetype(path, size)
                except Exception:
                    continue
        return ImageFont.load_default(size=size)

    # ── Shared helpers ────────────────────────────────────────────────────────
    def _norm(text: str) -> str:
        for src, dst in {
            "—": "-", "–": "-", "‒": "-",
            "‘": "'", "’": "'",
            "“": '"', "”": '"',
            "…": "...", " ": " ",
            "•": "*", "→": "->", "←": "<-",
            "©": "(c)", "®": "(R)", "™": "(TM)",
        }.items():
            text = text.replace(src, dst)
        return text

    def _wrap(text: str, font, max_px: int, draw) -> list[str]:
        if not text:
            return []
        words, lines, cur = text.split(), [], ""
        for word in words:
            test = (cur + " " + word).strip()
            if draw.textlength(test, font=font) <= max_px:
                cur = test
            else:
                if cur:
                    lines.append(cur)
                cur = word
        if cur:
            lines.append(cur)
        return lines

    def _block_h(lines, font, spacing, draw) -> int:
        total_h = 0
        for ln in lines:
            bb = draw.textbbox((0, 0), ln, font=font)
            total_h += (bb[3] - bb[1]) + spacing
        return total_h

    def _draw_block(draw, lines, font, x, y, fill, spacing=14,
                    centered=False, max_lines=99) -> int:
        for ln in lines[:max_lines]:
            bb = draw.textbbox((0, 0), ln, font=font)
            lw, lh = bb[2] - bb[0], bb[3] - bb[1]
            lx = (W - lw) // 2 if centered else x
            draw.text((lx, y), ln, font=font, fill=fill)
            y += lh + spacing
        return y

    def _rrect(draw, xy, radius, fill):
        draw.rounded_rectangle(list(xy), radius=radius, fill=fill)

    def _save(img: "Image.Image", idx: int) -> str:
        out = os.path.join(output_dir, f"slide_{idx:02d}.png")
        img.convert("RGB").save(out, "PNG", optimize=True)
        return out

    # ══════════════════════════════════════════════════════════════════════════
    # LAYOUT: listicle  (Bold Listicle)
    # White content slides, colored numbered badge circle top-left, navy footer
    # ══════════════════════════════════════════════════════════════════════════
    def _listicle_cover(idx, total, title, body) -> str:
        f_t = _load_font(76, bold=True)
        f_s = _load_font(40, bold=False)
        f_l = _load_font(26, bold=False)
        title, body = _norm(title), _norm(body)
        img = Image.new("RGB", (W, H), color=cover_bg)
        draw = ImageDraw.Draw(img, "RGBA")
        for row in range(H):
            draw.line([(0, row), (W, row)], fill=(*cover_accent, int(25 * (1 - row / H))))
        draw.rectangle([(0, 0), (10, H)], fill=cover_accent)
        pill = f"1 of {total}"
        pw = int(draw.textlength(pill, font=f_l)) + 36
        _rrect(draw, (50, 52, 50 + pw, 96), radius=22, fill=(*cover_accent, 200))
        draw.text((68, 60), pill, font=f_l, fill=cover_bg if sum(cover_accent) > 380 else WHITE)
        t_lines = _wrap(title, f_t, W - 140, draw)
        t_h = _block_h(t_lines[:4], f_t, 18, draw)
        t_y = max(150, (H // 2) - t_h // 2 - 80)
        y = _draw_block(draw, t_lines, f_t, 70, t_y, cover_text, spacing=18, centered=True, max_lines=4)
        draw.rectangle([(W // 2 - 50, y + 16), (W // 2 + 50, y + 22)], fill=cover_accent)
        s_lines = _wrap(body, f_s, W - 180, draw)
        _draw_block(draw, s_lines, f_s, 90, y + 46, fill=(*cover_text, 200), spacing=14, centered=True, max_lines=3)
        hint = "Swipe to read  >"
        hw = int(draw.textlength(hint, font=f_l))
        draw.text(((W - hw) // 2, H - 80), hint, font=f_l, fill=(*cover_accent, 160))
        return _save(img, idx)

    def _listicle_content(idx, total, title, body, badge_color) -> str:
        f_t = _load_font(64, bold=True)
        f_b = _load_font(38, bold=False)
        f_n = _load_font(50, bold=True)
        f_l = _load_font(26, bold=False)
        title, body = _norm(title), _norm(body)
        img = Image.new("RGB", (W, H), color=content_bg)
        draw = ImageDraw.Draw(img, "RGBA")
        draw.rectangle([(0, 0), (W, 10)], fill=badge_color)
        cx, cy, cr = 118, 155, 68
        draw.ellipse([(cx - cr, cy - cr), (cx + cr, cy + cr)], fill=badge_color)
        num_str = str(idx - 1)
        nw = int(draw.textlength(num_str, font=f_n))
        bb = draw.textbbox((0, 0), num_str, font=f_n)
        nh = bb[3] - bb[1]
        draw.text((cx - nw // 2, cy - nh // 2 - 3), num_str, font=f_n, fill=WHITE)
        PAD = 62
        t_lines = _wrap(title, f_t, W - PAD * 2, draw)
        y = _draw_block(draw, t_lines, f_t, PAD, cy + cr + 32, fill=title_color, spacing=12, max_lines=3)
        draw.rectangle([(PAD, y + 14), (PAD + 90, y + 20)], fill=badge_color)
        y += 48
        b_lines = _wrap(body, f_b, W - PAD * 2, draw)
        _draw_block(draw, b_lines, f_b, PAD, y, fill=body_color, spacing=18, max_lines=7)
        BAR = 80
        draw.rectangle([(0, H - BAR), (W, H)], fill=bottom_bar)
        draw.rectangle([(0, H - BAR), (int(W * idx / total), H - BAR + 5)], fill=badge_color)
        cnt = f"{idx} / {total}"
        cw = int(draw.textlength(cnt, font=f_l))
        draw.text((W - cw - 36, H - 52), cnt, font=f_l, fill=(*badge_color, 210))
        return _save(img, idx)

    def _listicle_cta(idx, total, title, body) -> str:
        f_t = _load_font(72, bold=True)
        f_s = _load_font(40, bold=False)
        f_l = _load_font(26, bold=False)
        title, body = _norm(title), _norm(body)
        img = Image.new("RGB", (W, H), color=cover_bg)
        draw = ImageDraw.Draw(img, "RGBA")
        for row in range(H):
            draw.line([(0, row), (W, row)], fill=(*cover_accent, int(25 * (row / H))))
        draw.rectangle([(W - 10, 0), (W, H)], fill=cover_accent)
        dx, dy, ds = 86, 96, 20
        draw.polygon([(dx, dy - ds), (dx + ds, dy), (dx, dy + ds), (dx - ds, dy)], fill=(*cover_accent, 180))
        pill = "Leave a comment below"
        pw = int(draw.textlength(pill, font=f_l)) + 36
        _rrect(draw, ((W - pw) // 2, 140, (W + pw) // 2, 184), radius=20, fill=(*cover_accent, 180))
        draw.text(((W - pw) // 2 + 18, 148), pill, font=f_l, fill=cover_bg if sum(cover_accent) > 380 else WHITE)
        cta_lines = _wrap(title, f_t, W - 140, draw)
        cta_h = _block_h(cta_lines[:3], f_t, 20, draw)
        cta_y = (H - cta_h) // 2 - 60
        y = _draw_block(draw, cta_lines, f_t, 70, cta_y, fill=cover_text, spacing=20, centered=True, max_lines=3)
        draw.rectangle([(W // 2 - 50, y + 18), (W // 2 + 50, y + 24)], fill=cover_accent)
        sub_lines = _wrap(body, f_s, W - 180, draw)
        _draw_block(draw, sub_lines, f_s, 90, y + 48, fill=(*cover_text, 190), spacing=16, centered=True, max_lines=3)
        draw.text((56, H - 76), f"{idx} / {total}", font=f_l, fill=(*cover_accent, 180))
        return _save(img, idx)

    # ══════════════════════════════════════════════════════════════════════════
    # LAYOUT: dark_minimal  (Minimal Dark)
    # Near-black slides; HUGE left-aligned title; thin gold rule; no badge circle
    # Left border bar is the only accent element on content slides
    # ══════════════════════════════════════════════════════════════════════════
    def _dark_cover(idx, total, title, body) -> str:
        f_t = _load_font(82, bold=True)
        f_s = _load_font(36, bold=False)
        f_l = _load_font(24, bold=False)
        title, body = _norm(title), _norm(body)
        img = Image.new("RGB", (W, H), color=cover_bg)
        draw = ImageDraw.Draw(img, "RGBA")
        # Large decorative background character
        f_bg = _load_font(500, bold=True)
        draw.text((-40, H // 2 - 280), '"', font=f_bg, fill=(*cover_accent, 15))
        # Left thick accent bar
        draw.rectangle([(0, 0), (8, H)], fill=cover_accent)
        # Top-right slide counter
        cnt = f"01 / {total:02d}"
        cw = int(draw.textlength(cnt, font=f_l))
        draw.text((W - cw - 50, 52), cnt, font=f_l, fill=(*cover_accent, 160))
        # Title — large, left-aligned, starts at 30% from top
        PAD = 70
        t_lines = _wrap(title, f_t, W - PAD - 60, draw)
        y = _draw_block(draw, t_lines, f_t, PAD, max(220, H // 3 - 60), cover_text, spacing=14, max_lines=4)
        # Thin gold rule
        draw.rectangle([(PAD, y + 24), (PAD + 120, y + 28)], fill=cover_accent)
        y += 52
        s_lines = _wrap(body, f_s, W - PAD - 60, draw)
        _draw_block(draw, s_lines, f_s, PAD, y, fill=(*cover_text, 170), spacing=16, max_lines=4)
        # Bottom "Swipe" hint
        hint = "SWIPE  >"
        draw.text((PAD, H - 80), hint, font=f_l, fill=(*cover_accent, 120))
        return _save(img, idx)

    def _dark_content(idx, total, title, body, badge_color) -> str:
        f_t = _load_font(70, bold=True)
        f_b = _load_font(36, bold=False)
        f_l = _load_font(24, bold=False)
        title, body = _norm(title), _norm(body)
        img = Image.new("RGB", (W, H), color=content_bg)
        draw = ImageDraw.Draw(img, "RGBA")
        # Left accent bar (thicker than cover)
        draw.rectangle([(0, 0), (16, H)], fill=badge_color)
        # Slide number — top right, muted
        cnt_str = f"{idx:02d} / {total:02d}"
        cw = int(draw.textlength(cnt_str, font=f_l))
        draw.text((W - cw - 50, 52), cnt_str, font=f_l, fill=(*badge_color, 140))
        # LARGE title left-aligned, starts high
        PAD = 70
        t_lines = _wrap(title, f_t, W - PAD - 60, draw)
        y = _draw_block(draw, t_lines, f_t, PAD, 180, title_color, spacing=14, max_lines=4)
        # Thin colored rule
        draw.rectangle([(PAD, y + 24), (PAD + 100, y + 27)], fill=badge_color)
        y += 52
        # Body text — smaller, muted
        b_lines = _wrap(body, f_b, W - PAD - 60, draw)
        _draw_block(draw, b_lines, f_b, PAD, y, fill=body_color, spacing=20, max_lines=8)
        # Bottom: thin line only
        draw.rectangle([(0, H - 60), (W, H - 58)], fill=(*badge_color, 80))
        return _save(img, idx)

    def _dark_cta(idx, total, title, body) -> str:
        f_t = _load_font(78, bold=True)
        f_s = _load_font(34, bold=False)
        f_l = _load_font(24, bold=False)
        title, body = _norm(title), _norm(body)
        img = Image.new("RGB", (W, H), color=cover_bg)
        draw = ImageDraw.Draw(img, "RGBA")
        draw.rectangle([(0, 0), (8, H)], fill=cover_accent)
        cnt_str = f"{idx:02d} / {total:02d}"
        cw = int(draw.textlength(cnt_str, font=f_l))
        draw.text((W - cw - 50, 52), cnt_str, font=f_l, fill=(*cover_accent, 140))
        PAD = 70
        # CTA title — center of slide
        t_lines = _wrap(title, f_t, W - PAD - 60, draw)
        t_h = _block_h(t_lines[:3], f_t, 16, draw)
        t_y = (H - t_h) // 2 - 80
        y = _draw_block(draw, t_lines, f_t, PAD, t_y, cover_text, spacing=16, max_lines=3)
        draw.rectangle([(PAD, y + 24), (PAD + 100, y + 27)], fill=cover_accent)
        y += 52
        s_lines = _wrap(body, f_s, W - PAD - 60, draw)
        _draw_block(draw, s_lines, f_s, PAD, y, fill=(*cover_text, 170), spacing=18, max_lines=4)
        return _save(img, idx)

    # ══════════════════════════════════════════════════════════════════════════
    # LAYOUT: stat_big  (Stat Reveal)
    # Blue cover; content: title displayed ENORMOUS centered; tiny body below
    # ══════════════════════════════════════════════════════════════════════════
    def _stat_cover(idx, total, title, body) -> str:
        f_t = _load_font(74, bold=True)
        f_s = _load_font(38, bold=False)
        f_l = _load_font(24, bold=False)
        title, body = _norm(title), _norm(body)
        img = Image.new("RGB", (W, H), color=cover_bg)
        draw = ImageDraw.Draw(img, "RGBA")
        # Wave-like two-tone split
        draw.rectangle([(0, H - 220), (W, H)], fill=(*cover_accent, 30))
        # Decorative large number "?" hinting at reveals
        f_deco = _load_font(300, bold=True)
        draw.text((W - 220, H // 2 - 200), "?", font=f_deco, fill=(*cover_accent, 18))
        # Slide counter pill top-center
        pill = f"1 of {total} reveals"
        pw = int(draw.textlength(pill, font=f_l)) + 36
        _rrect(draw, ((W - pw) // 2, 52, (W + pw) // 2, 96), radius=22, fill=(*cover_accent, 200))
        draw.text(((W - pw) // 2 + 18, 60), pill, font=f_l, fill=cover_bg)
        # Title centered
        PAD = 60
        t_lines = _wrap(title, f_t, W - PAD * 2, draw)
        t_h = _block_h(t_lines[:4], f_t, 18, draw)
        y = _draw_block(draw, t_lines, f_t, PAD, max(180, (H // 2) - t_h // 2 - 60), cover_text,
                        spacing=18, centered=True, max_lines=4)
        draw.rectangle([(W // 2 - 40, y + 20), (W // 2 + 40, y + 24)], fill=cover_accent)
        y += 48
        s_lines = _wrap(body, f_s, W - PAD * 2, draw)
        _draw_block(draw, s_lines, f_s, PAD, y, fill=(*cover_text, 200), spacing=14, centered=True, max_lines=3)
        return _save(img, idx)

    def _stat_content(idx, total, title, body, badge_color) -> str:
        f_huge = _load_font(100, bold=True)  # title as massive centered text
        f_body = _load_font(36, bold=False)
        f_l    = _load_font(24, bold=False)
        f_num  = _load_font(22, bold=True)
        title, body = _norm(title), _norm(body)
        img = Image.new("RGB", (W, H), color=content_bg)
        draw = ImageDraw.Draw(img, "RGBA")
        # Top color band
        draw.rectangle([(0, 0), (W, 90)], fill=badge_color)
        # Step number in top band
        step_label = f"#{idx - 1}"
        draw.text((W // 2 - int(draw.textlength(step_label, font=f_num)) // 2, 30),
                  step_label, font=f_num, fill=(*content_bg, 220))
        PAD = 60
        # HUGE title — centered vertically in upper 65%
        t_lines = _wrap(title, f_huge, W - PAD * 2, draw)
        t_h = _block_h(t_lines[:3], f_huge, 18, draw)
        t_y = max(130, (H * 65 // 100) // 2 - t_h // 2)
        y = _draw_block(draw, t_lines, f_huge, PAD, t_y, title_color, spacing=18, centered=True, max_lines=3)
        # Horizontal rule
        draw.rectangle([(PAD, y + 22), (W - PAD, y + 25)], fill=(*badge_color, 120))
        y += 50
        # Body small, centered
        b_lines = _wrap(body, f_body, W - PAD * 2, draw)
        _draw_block(draw, b_lines, f_body, PAD, y, fill=body_color, spacing=18, centered=True, max_lines=5)
        # Bottom strip
        draw.rectangle([(0, H - 60), (W, H)], fill=bottom_bar)
        progress_w = int(W * (idx - 1) / max(total - 2, 1))
        draw.rectangle([(0, H - 60), (progress_w, H - 55)], fill=badge_color)
        cnt = f"{idx - 1} / {total - 2}"
        cw = int(draw.textlength(cnt, font=f_l))
        draw.text((W - cw - 36, H - 42), cnt, font=f_l, fill=(*badge_color, 200))
        return _save(img, idx)

    def _stat_cta(idx, total, title, body) -> str:
        f_t = _load_font(74, bold=True)
        f_s = _load_font(38, bold=False)
        f_l = _load_font(24, bold=False)
        title, body = _norm(title), _norm(body)
        img = Image.new("RGB", (W, H), color=cover_bg)
        draw = ImageDraw.Draw(img, "RGBA")
        draw.rectangle([(0, 0), (W, H)], fill=cover_bg)
        # Decorative corner triangle
        draw.polygon([(0, 0), (300, 0), (0, 300)], fill=(*cover_accent, 40))
        draw.polygon([(W, H), (W - 300, H), (W, H - 300)], fill=(*cover_accent, 40))
        PAD = 70
        t_lines = _wrap(title, f_t, W - PAD * 2, draw)
        t_h = _block_h(t_lines[:3], f_t, 18, draw)
        y = _draw_block(draw, t_lines, f_t, PAD, (H - t_h) // 2 - 80,
                        cover_text, spacing=18, centered=True, max_lines=3)
        draw.rectangle([(W // 2 - 40, y + 22), (W // 2 + 40, y + 25)], fill=cover_accent)
        y += 50
        s_lines = _wrap(body, f_s, W - PAD * 2, draw)
        _draw_block(draw, s_lines, f_s, PAD, y, fill=(*cover_text, 190), spacing=16, centered=True, max_lines=3)
        draw.text((W // 2 - int(draw.textlength(f"{idx} / {total}", font=f_l)) // 2, H - 70),
                  f"{idx} / {total}", font=f_l, fill=(*cover_accent, 160))
        return _save(img, idx)

    # ══════════════════════════════════════════════════════════════════════════
    # LAYOUT: step_progress  (Step Framework)
    # Green cover; content: top step-indicator strip, arrow-bulleted body
    # ══════════════════════════════════════════════════════════════════════════
    def _step_cover(idx, total, title, body) -> str:
        f_t = _load_font(76, bold=True)
        f_s = _load_font(38, bold=False)
        f_l = _load_font(24, bold=False)
        title, body = _norm(title), _norm(body)
        img = Image.new("RGB", (W, H), color=cover_bg)
        draw = ImageDraw.Draw(img, "RGBA")
        # Diagonal accent block bottom-right
        draw.polygon([(W - 280, H), (W, H - 280), (W, H)], fill=(*cover_accent, 60))
        # Step label strip at top
        draw.rectangle([(0, 0), (W, 100)], fill=(*WHITE, 20))
        header = f"A {total - 2}-Step Framework"
        hw = int(draw.textlength(header, font=f_l))
        draw.text(((W - hw) // 2, 36), header, font=f_l, fill=(*cover_accent, 240))
        # Title centered
        PAD = 70
        t_lines = _wrap(title, f_t, W - PAD * 2, draw)
        t_h = _block_h(t_lines[:4], f_t, 18, draw)
        y = _draw_block(draw, t_lines, f_t, PAD, max(170, (H // 2) - t_h // 2 - 40),
                        cover_text, spacing=18, centered=True, max_lines=4)
        draw.rectangle([(PAD, y + 20), (W - PAD, y + 24)], fill=(*cover_accent, 180))
        y += 48
        s_lines = _wrap(body, f_s, W - PAD * 2, draw)
        _draw_block(draw, s_lines, f_s, PAD, y, fill=(*cover_text, 200), spacing=14, centered=True, max_lines=3)
        return _save(img, idx)

    def _step_content(idx, total, title, body, badge_color) -> str:
        f_t = _load_font(62, bold=True)
        f_b = _load_font(36, bold=False)
        f_n = _load_font(30, bold=True)
        title, body = _norm(title), _norm(body)
        content_steps = total - 2  # exclude cover + CTA
        step_num = max(1, idx - 1)

        img = Image.new("RGB", (W, H), color=content_bg)
        draw = ImageDraw.Draw(img, "RGBA")

        # ── Step progress strip at top ────────────────────────────────────────
        STRIP = 100
        draw.rectangle([(0, 0), (W, STRIP)], fill=bottom_bar)
        dot_r = 18
        dot_gap = max(50, min(90, (W - 140) // max(content_steps, 1)))
        start_x = (W - (content_steps * dot_gap - (dot_gap - dot_r * 2))) // 2

        for i in range(content_steps):
            cx = start_x + i * dot_gap + dot_r
            cy = STRIP // 2
            if i + 1 < step_num:
                # completed — filled
                draw.ellipse([(cx - dot_r, cy - dot_r), (cx + dot_r, cy + dot_r)], fill=badge_color)
                check = "+"
                cw = int(draw.textlength(check, font=f_n))
                draw.text((cx - cw // 2, cy - 14), check, font=f_n, fill=WHITE)
            elif i + 1 == step_num:
                # current — filled with border
                draw.ellipse([(cx - dot_r - 4, cy - dot_r - 4), (cx + dot_r + 4, cy + dot_r + 4)],
                             fill=WHITE)
                draw.ellipse([(cx - dot_r, cy - dot_r), (cx + dot_r, cy + dot_r)], fill=badge_color)
                num = str(step_num)
                nw = int(draw.textlength(num, font=f_n))
                draw.text((cx - nw // 2, cy - 14), num, font=f_n, fill=WHITE)
            else:
                # future — outline only
                draw.ellipse([(cx - dot_r, cy - dot_r), (cx + dot_r, cy + dot_r)],
                             outline=(*badge_color, 100), width=3)
            # connecting line between dots
            if i < content_steps - 1:
                nx = start_x + (i + 1) * dot_gap + dot_r
                line_color = badge_color if i + 1 < step_num else (*badge_color, 60)
                draw.line([(cx + dot_r, cy), (nx - dot_r, cy)], fill=line_color, width=3)

        # ── Title ─────────────────────────────────────────────────────────────
        PAD = 62
        t_lines = _wrap(title, f_t, W - PAD * 2, draw)
        y = _draw_block(draw, t_lines, f_t, PAD, STRIP + 40, title_color, spacing=12, max_lines=3)

        # Accent underline
        draw.rectangle([(PAD, y + 12), (PAD + 80, y + 17)], fill=badge_color)
        y += 42

        # ── Body with arrow bullets ───────────────────────────────────────────
        for line_text in _wrap(body, f_b, W - PAD - 30, draw)[:7]:
            draw.text((PAD, y), "->", font=f_b, fill=badge_color)
            draw.text((PAD + 52, y), line_text, font=f_b, fill=body_color)
            bb = draw.textbbox((0, 0), line_text, font=f_b)
            y += (bb[3] - bb[1]) + 20

        # ── Bottom bar ────────────────────────────────────────────────────────
        BAR = 60
        draw.rectangle([(0, H - BAR), (W, H)], fill=bottom_bar)
        return _save(img, idx)

    def _step_cta(idx, total, title, body) -> str:
        f_t = _load_font(74, bold=True)
        f_s = _load_font(36, bold=False)
        title, body = _norm(title), _norm(body)
        img = Image.new("RGB", (W, H), color=cover_bg)
        draw = ImageDraw.Draw(img, "RGBA")
        draw.polygon([(W - 280, H), (W, H - 280), (W, H)], fill=(*cover_accent, 60))
        draw.polygon([(0, 0), (280, 0), (0, 280)], fill=(*cover_accent, 40))
        # Checkmark large
        f_check = _load_font(120, bold=True)
        check_text = "Done!"
        cw = int(draw.textlength(check_text, font=f_check))
        draw.text(((W - cw) // 2, 120), check_text, font=f_check, fill=(*cover_accent, 220))
        PAD = 70
        t_lines = _wrap(title, f_t, W - PAD * 2, draw)
        y = _draw_block(draw, t_lines, f_t, PAD, 380, cover_text, spacing=18, centered=True, max_lines=3)
        draw.rectangle([(PAD, y + 20), (W - PAD, y + 23)], fill=(*cover_accent, 180))
        y += 48
        s_lines = _wrap(body, f_s, W - PAD * 2, draw)
        _draw_block(draw, s_lines, f_s, PAD, y, fill=(*cover_text, 200), spacing=14, centered=True, max_lines=3)
        return _save(img, idx)

    # ══════════════════════════════════════════════════════════════════════════
    # LAYOUT: quote_pull  (Story Arc)
    # Cream content slides; giant quotation mark; editorial/magazine feel
    # Right amber border bar; title reads like a pull-quote
    # ══════════════════════════════════════════════════════════════════════════
    def _story_cover(idx, total, title, body) -> str:
        f_t = _load_font(74, bold=True)
        f_s = _load_font(38, bold=False)
        f_l = _load_font(24, bold=False)
        f_big = _load_font(260, bold=True)
        title, body = _norm(title), _norm(body)
        img = Image.new("RGB", (W, H), color=cover_bg)
        draw = ImageDraw.Draw(img, "RGBA")
        # Big decorative quote mark — watermark
        draw.text((50, 100), '"', font=f_big, fill=(*cover_accent, 35))
        # Warm texture: bottom band
        draw.rectangle([(0, H - 180), (W, H)], fill=(*cover_accent, 40))
        # Right border
        draw.rectangle([(W - 12, 0), (W, H)], fill=cover_accent)
        # Issue/series label top
        label = f"Part 1 of {total}"
        lw = int(draw.textlength(label, font=f_l))
        draw.text((W - lw - 36, 48), label, font=f_l, fill=(*cover_accent, 200))
        PAD = 70
        t_lines = _wrap(title, f_t, W - PAD - 80, draw)
        t_h = _block_h(t_lines[:4], f_t, 16, draw)
        y = _draw_block(draw, t_lines, f_t, PAD, max(200, (H // 2) - t_h // 2 - 60),
                        cover_text, spacing=16, max_lines=4)
        # Amber rule
        draw.rectangle([(PAD, y + 20), (PAD + 100, y + 24)], fill=cover_accent)
        y += 50
        s_lines = _wrap(body, f_s, W - PAD - 80, draw)
        _draw_block(draw, s_lines, f_s, PAD, y, fill=(*cover_text, 190), spacing=16, max_lines=3)
        return _save(img, idx)

    def _story_content(idx, total, title, body, badge_color) -> str:
        f_quote = _load_font(180, bold=True)
        f_t     = _load_font(58, bold=True)
        f_b     = _load_font(36, bold=False)
        f_l     = _load_font(22, bold=False)
        title, body = _norm(title), _norm(body)
        img = Image.new("RGB", (W, H), color=content_bg)
        draw = ImageDraw.Draw(img, "RGBA")
        # Right accent border
        draw.rectangle([(W - 14, 0), (W, H)], fill=badge_color)
        # Slide number — top-left badge (square, not circle)
        draw.rectangle([(50, 50), (116, 116)], fill=badge_color)
        num_str = str(idx - 1)
        f_n = _load_font(42, bold=True)
        nw = int(draw.textlength(num_str, font=f_n))
        bb = draw.textbbox((0, 0), num_str, font=f_n)
        nh = bb[3] - bb[1]
        draw.text((83 - nw // 2, 83 - nh // 2 - 3), num_str, font=f_n, fill=WHITE)
        # Large decorative quote mark
        draw.text((50, 90), '"', font=f_quote, fill=(*badge_color, 25))
        # Title — larger, treated as pull-quote
        PAD = 70
        t_lines = _wrap(title, f_t, W - PAD - 80, draw)
        y = _draw_block(draw, t_lines, f_t, PAD, 220, title_color, spacing=14, max_lines=4)
        # Amber rule
        draw.rectangle([(PAD, y + 16), (PAD + 100, y + 20)], fill=badge_color)
        y += 46
        # Body text — softer
        b_lines = _wrap(body, f_b, W - PAD - 80, draw)
        _draw_block(draw, b_lines, f_b, PAD, y, fill=body_color, spacing=20, max_lines=6)
        # Bottom: issue label
        draw.rectangle([(0, H - 60), (W - 14, H)], fill=bottom_bar)
        cnt = f"Part {idx - 1} of {total - 2}"
        draw.text((PAD, H - 44), cnt, font=f_l, fill=(*badge_color, 200))
        return _save(img, idx)

    def _story_cta(idx, total, title, body) -> str:
        f_quote = _load_font(180, bold=True)
        f_t = _load_font(68, bold=True)
        f_s = _load_font(36, bold=False)
        f_l = _load_font(24, bold=False)
        title, body = _norm(title), _norm(body)
        img = Image.new("RGB", (W, H), color=cover_bg)
        draw = ImageDraw.Draw(img, "RGBA")
        # Closing quote mark (right-aligned)
        draw.text((W - 200, H // 2 - 80), '"', font=f_quote, fill=(*cover_accent, 30))
        draw.rectangle([(W - 12, 0), (W, H)], fill=cover_accent)
        # "The End" style label
        label = "The Takeaway"
        lw = int(draw.textlength(label, font=f_l))
        draw.text((W - lw - 36, 48), label, font=f_l, fill=(*cover_accent, 200))
        PAD = 70
        t_lines = _wrap(title, f_t, W - PAD - 80, draw)
        t_h = _block_h(t_lines[:3], f_t, 16, draw)
        y = _draw_block(draw, t_lines, f_t, PAD, max(200, (H // 2) - t_h // 2 - 80),
                        cover_text, spacing=16, max_lines=3)
        draw.rectangle([(PAD, y + 20), (PAD + 100, y + 24)], fill=cover_accent)
        y += 50
        s_lines = _wrap(body, f_s, W - PAD - 80, draw)
        _draw_block(draw, s_lines, f_s, PAD, y, fill=(*cover_text, 190), spacing=16, max_lines=3)
        draw.text((PAD, H - 70), f"Part {idx} of {total}", font=f_l, fill=(*cover_accent, 160))
        return _save(img, idx)

    # ══════════════════════════════════════════════════════════════════════════
    # Dispatch to the right render functions
    # ══════════════════════════════════════════════════════════════════════════
    COVER_FN = {
        "listicle":      _listicle_cover,
        "dark_minimal":  _dark_cover,
        "stat_big":      _stat_cover,
        "step_progress": _step_cover,
        "quote_pull":    _story_cover,
    }
    CONTENT_FN = {
        "listicle":      _listicle_content,
        "dark_minimal":  _dark_content,
        "stat_big":      _stat_content,
        "step_progress": _step_content,
        "quote_pull":    _story_content,
    }
    CTA_FN = {
        "listicle":      _listicle_cta,
        "dark_minimal":  _dark_cta,
        "stat_big":      _stat_cta,
        "step_progress": _step_cta,
        "quote_pull":    _story_cta,
    }
    render_cover   = COVER_FN.get(layout, _listicle_cover)
    render_content = CONTENT_FN.get(layout, _listicle_content)
    render_cta     = CTA_FN.get(layout, _listicle_cta)

    # ── Collect slides from carousel model ────────────────────────────────────
    slides_data: list[tuple[str, str]] = []

    def _add(t, c):
        slides_data.append((t or "", c or ""))

    if isinstance(carousel_data, EducationalContentCarousel):
        _add(carousel_data.cover.title, carousel_data.cover.content)
        for s in carousel_data.contents:
            _add(s.title, s.content)
        _add(carousel_data.call_to_action.title, carousel_data.call_to_action.content)
    elif isinstance(carousel_data, CaseStudyCarousel):
        _add(carousel_data.cover.title, carousel_data.cover.content)
        _add(carousel_data.challenge.title, carousel_data.challenge.content)
        _add(carousel_data.solution.title, carousel_data.solution.content)
        _add(carousel_data.results.title, carousel_data.results.content)
        if carousel_data.testimonial:
            _add(carousel_data.testimonial.title, carousel_data.testimonial.content)
        _add(carousel_data.call_to_action.title, carousel_data.call_to_action.content)
    elif isinstance(carousel_data, PersonalStoryCarousel):
        _add(carousel_data.cover.title, carousel_data.cover.content)
        for s in carousel_data.story_slides:
            _add(s.title, s.content)
        _add(carousel_data.takeaway.title, carousel_data.takeaway.content)
        _add(carousel_data.call_to_action.title, carousel_data.call_to_action.content)
    elif isinstance(carousel_data, IndustryInsightsCarousel):
        _add(carousel_data.cover.title, carousel_data.cover.content)
        for s in carousel_data.insights:
            _add(s.title, s.content)
        _add(carousel_data.call_to_action.title, carousel_data.call_to_action.content)
    elif isinstance(carousel_data, EventRecapCarousel):
        _add(carousel_data.cover.title, carousel_data.cover.content)
        for s in carousel_data.key_moments:
            _add(s.title, s.content)
        _add(carousel_data.call_to_action.title, carousel_data.call_to_action.content)
    elif isinstance(carousel_data, TestimonialCarousel):
        _add(carousel_data.cover.title, carousel_data.cover.content)
        for s in carousel_data.testimonials:
            _add(s.title, s.content)
        _add(carousel_data.call_to_action.title, carousel_data.call_to_action.content)
    elif isinstance(carousel_data, ProductDemoCarousel):
        _add(carousel_data.cover.title, carousel_data.cover.content)
        _add(carousel_data.main_feature.title, carousel_data.main_feature.content)
        for s in carousel_data.additional_features:
            _add(s.title, s.content)
        _add(carousel_data.call_to_action.title, carousel_data.call_to_action.content)

    # ── Render ────────────────────────────────────────────────────────────────
    total = len(slides_data)
    image_paths = []
    for idx, (title, body) in enumerate(slides_data, start=1):
        if idx == 1:
            path = render_cover(idx, total, title, body)
        elif idx == total:
            path = render_cta(idx, total, title, body)
        else:
            bc = badge_colors[(idx - 2) % len(badge_colors)]
            path = render_content(idx, total, title, body, bc)
        image_paths.append(path)

    return image_paths


if __name__ == "__main__":
    # Debug
    # debug_master_slide_placeholders_and_text()

    # test_caption_only_slide()
    # test_create_educational_ppt()
    test_create_case_study_ppt()

    exit(0)

    # Example usage of the TestimonialCarousel model
    carousel_data = {
        "cover": {
            "title": "What Our Clients Are Saying",
            "content": "Hear from our satisfied clients about their experience "
        },
        "testimonials": [
            {
                "title": "Client Testimonial 1",
                "content": "Lorem ipsum dolor sit amet, consectetur adipiscing elit. Sed do eiusmod tempor incididunt ut labore et dolore magna aliqua.",
                "image_url": "https://example.com/client1.jpg",
                "client_name": "John Doe",
                "client_logo_url": "https://example.com/logo1.jpg"
            },
            {
                "title": "Client Testimonial 2",
                "content": "Ut enim ad minim veniam, quis nostrud exercitation ullamco laboris nisi ut aliquip ex ea commodo consequat.",
                "image_url": "https://example.com/client2.jpg",
                "client_name": "Jane Smith",
                "client_logo_url": "https://example.com/logo2.jpg"
            }
        ],
        "call_to_action": {
            "title": "Contact Us",
            "content": "If you are ready tto learn more about our services."
        }

    }
    ppt = create_ppt("Testimonials", TestimonialCarousel(**carousel_data))
    print(f"Presentation created: {ppt}")
